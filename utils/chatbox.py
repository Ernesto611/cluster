from chatbot.models import ChatbotDocumento
import fitz
import docx
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from django.core.exceptions import ValidationError
import os
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def cargar_documentos_chatbox():

    contenido = []
    documentos = ChatbotDocumento.objects.filter(activo=True)

    for doc in documentos:
        try:
            filename = doc.archivo.name.lower()
            path = doc.archivo.path

            if not os.path.exists(path):
                logger.warning(f"⚠️ Archivo no encontrado: {doc.nombre}")
                continue

            file_size = os.path.getsize(path)
            if file_size > 50 * 1024 * 1024:
                logger.warning(f"⚠️ Archivo demasiado grande: {doc.nombre} ({file_size / 1024 / 1024:.1f}MB)")
                continue

            if filename.endswith(".pdf"):
                texto = extraer_texto_pdf(path)
            elif filename.endswith(".docx"):
                texto = extraer_texto_docx(path)
            elif filename.endswith(".pptx"):
                texto = extraer_texto_pptx(path)
            elif filename.endswith(".xlsx"):
                texto = extraer_texto_xlsx(path)
            elif filename.endswith((".jpg", ".jpeg", ".png", ".webp")):
                texto = extraer_texto_imagen(path)
            else:
                logger.warning(f"⚠️ Formato no soportado: {doc.nombre}")
                continue

            if texto and len(texto.strip()) > 0:
                contenido.append(texto)
            else:
                logger.warning(f"⚠️ No se pudo extraer texto de: {doc.nombre}")

        except Exception as e:
            logger.error(f"⚠️ Error leyendo {doc.nombre}: {e}")
            continue

    return "\n".join(contenido)

def extraer_texto_pdf(filepath):

    texto = ""
    try:
        with fitz.open(filepath) as pdf:
            for page_num, page in enumerate(pdf):
                page_text = page.get_text()

                if len(page_text.strip()) < 50:
                    logger.info(f"Página {page_num + 1} parece ser imagen escaneada, usando OCR...")
                    try:

                        mat = fitz.Matrix(2.0, 2.0)
                        pix = page.get_pixmap(matrix=mat)
                        img_data = pix.tobytes("png")

                        from io import BytesIO
                        image = Image.open(BytesIO(img_data))

                        ocr_text = pytesseract.image_to_string(image, lang="spa")
                        if len(ocr_text.strip()) > len(page_text.strip()):
                            page_text = ocr_text

                    except Exception as ocr_error:
                        logger.warning(f"Error en OCR para página {page_num + 1}: {ocr_error}")

                texto += page_text + "\n"

        if not texto.strip():
            raise ValueError("No se pudo extraer texto del PDF")

    except Exception as e:
        logger.error(f"Error procesando PDF {filepath}: {e}")
        raise

    return texto.strip()

def extraer_texto_docx(filepath):

    texto = ""
    texto_ocr = ""

    try:
        doc = docx.Document(filepath)

        for para in doc.paragraphs:
            texto += para.text + "\n"

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    texto += cell.text + "\t"
                texto += "\n"

        if len(texto.strip()) < 100:
            logger.info("Documento DOCX con poco texto, buscando imágenes para OCR...")
            texto_ocr = extraer_imagenes_docx_ocr(doc)

            if texto_ocr:
                texto += "\n--- TEXTO EXTRAÍDO DE IMÁGENES ---\n" + texto_ocr
                logger.info(f"OCR extrajo {len(texto_ocr)} caracteres adicionales de imágenes")

        if not texto.strip():
            raise ValueError("No se pudo extraer texto del archivo DOCX")

    except Exception as e:
        logger.error(f"Error procesando DOCX {filepath}: {e}")
        raise

    return texto.strip()

def extraer_imagenes_docx_ocr(doc):

    from docx.oxml.ns import nsdecls, qn
    from docx.oxml.parser import parse_xml
    import tempfile
    import os

    texto_ocr = ""

    try:

        image_parts = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_parts.append(rel.target_part)

        for i, image_part in enumerate(image_parts):
            try:

                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                    tmp_file.write(image_part.blob)
                    tmp_path = tmp_file.name

                try:
                    with Image.open(tmp_path) as image:

                        if image.mode != 'RGB':
                            image = image.convert('RGB')

                        width, height = image.size
                        if width < 800 or height < 600:
                            new_width = max(width * 2, 800)
                            new_height = max(height * 2, 600)
                            image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)

                        custom_config = r'--oem 3 --psm 6'
                        ocr_text = pytesseract.image_to_string(image, lang="spa", config=custom_config)

                        if len(ocr_text.strip()) > 10:
                            texto_ocr += f"\n--- Imagen {i+1} ---\n{ocr_text.strip()}\n"

                except Exception as ocr_error:
                    logger.warning(f"Error aplicando OCR a imagen {i+1} en DOCX: {ocr_error}")

                os.unlink(tmp_path)

            except Exception as img_error:
                logger.warning(f"Error procesando imagen {i+1} en DOCX: {img_error}")
                continue

    except Exception as e:
        logger.warning(f"Error general extrayendo imágenes de DOCX: {e}")

    return texto_ocr

def extraer_texto_pptx(filepath):

    texto = ""
    try:
        prs = Presentation(filepath)

        for slide_num, slide in enumerate(prs.slides):
            texto += f"--- Diapositiva {slide_num + 1} ---\n"
            slide_text = ""

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"

            if len(slide_text.strip()) < 50:
                logger.info(f"Diapositiva {slide_num + 1} con poco texto, buscando imágenes para OCR...")
                ocr_text = extraer_imagenes_pptx_ocr(slide, slide_num + 1)
                if ocr_text:
                    slide_text += f"\n--- TEXTO DE IMÁGENES (Diapositiva {slide_num + 1}) ---\n{ocr_text}"

            texto += slide_text + "\n"

        if not texto.strip():
            raise ValueError("No se pudo extraer texto del archivo PPTX")

    except Exception as e:
        logger.error(f"Error procesando PPTX {filepath}: {e}")
        raise

    return texto.strip()

def extraer_imagenes_pptx_ocr(slide, slide_num):

    import tempfile
    import os
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    texto_ocr = ""

    try:
        imagen_count = 0

        for shape in slide.shapes:

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    imagen_count += 1

                    image_part = shape.part.related_parts[shape._element.blip_rId]

                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                        tmp_file.write(image_part.blob)
                        tmp_path = tmp_file.name

                    try:
                        with Image.open(tmp_path) as image:

                            if image.mode not in ['RGB', 'L']:
                                image = image.convert('RGB')

                            width, height = image.size
                            if width < 800 or height < 600:
                                scale_factor = max(800/width, 600/height, 2.0)
                                new_width = int(width * scale_factor)
                                new_height = int(height * scale_factor)
                                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)

                            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789áéíóúüñÁÉÍÓÚÜÑ.,;:!?¿¡()[]{}"-+*/=@#$%&_ '
                            ocr_text = pytesseract.image_to_string(image, lang="spa", config=custom_config)

                            if len(ocr_text.strip()) > 5:

                                ocr_text_clean = limpiar_texto_ocr(ocr_text)
                                if len(ocr_text_clean.strip()) > 5:
                                    texto_ocr += f"Imagen {imagen_count}: {ocr_text_clean.strip()}\n"

                    except Exception as ocr_error:
                        logger.warning(f"Error OCR en imagen {imagen_count} de diapositiva {slide_num}: {ocr_error}")

                    os.unlink(tmp_path)

                except Exception as img_error:
                    logger.warning(f"Error procesando imagen {imagen_count} en diapositiva {slide_num}: {img_error}")
                    continue

            elif hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                try:

                    if shape.fill.type == 6:

                        pass
                except:
                    pass

    except Exception as e:
        logger.warning(f"Error general extrayendo imágenes de diapositiva {slide_num}: {e}")

    return texto_ocr

def limpiar_texto_ocr(texto):

    import re

    texto = re.sub(r'[^\w\sáéíóúüñÁÉÍÓÚÜÑ.,;:!?¿¡()[\]{}"\-+*/=@#$%&_]', '', texto)

    texto = re.sub(r'\s+', ' ', texto)

    lineas = texto.split('\n')
    lineas_limpias = [linea.strip() for linea in lineas if len(linea.strip()) > 2]

    return '\n'.join(lineas_limpias)

def extraer_texto_xlsx(filepath):

    texto = ""
    try:
        wb = load_workbook(filename=filepath, read_only=True, data_only=True)

        for sheet in wb.worksheets:
            texto += f"=== Hoja: {sheet.title} ===\n"

            non_empty_cells = 0
            sheet_content = ""

            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:

                        if isinstance(cell, (int, float)):
                            if cell == int(cell):
                                row_text.append(str(int(cell)))
                            else:
                                row_text.append(f"{cell:.2f}")
                        else:
                            row_text.append(str(cell))
                        non_empty_cells += 1

                if row_text:
                    sheet_content += "\t".join(row_text) + "\n"

            if non_empty_cells < 5:
                logger.warning(f"Hoja '{sheet.title}' tiene muy pocas celdas con datos ({non_empty_cells})")

            if any(char in sheet_content for char in ['�', '□', '▯']):
                logger.warning(f"Hoja '{sheet.title}' contiene símbolos o caracteres no reconocidos")
                sheet_content += "\n[ADVERTENCIA: Contenido con símbolos no reconocidos detectado]\n"

            texto += sheet_content + "\n"

        if not texto.strip():
            raise ValueError("No se pudo extraer texto del archivo Excel")

    except Exception as e:
        logger.error(f"Error procesando Excel {filepath}: {e}")
        raise

    return texto.strip()

def extraer_texto_imagen(filepath):

    try:

        with Image.open(filepath) as image:

            width, height = image.size
            if width < 100 or height < 100:
                logger.warning(f"Imagen muy pequeña ({width}x{height}), el OCR podría ser impreciso")

            if image.mode != 'RGB':
                image = image.convert('RGB')

            if width < 500 or height < 500:

                new_width = max(width * 2, 800)
                new_height = max(height * 2, 600)
                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
                logger.info("Imagen redimensionada para mejorar OCR")

            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789áéíóúüñÁÉÍÓÚÜÑ.,;:!?¿¡()[]{}"-+*/=@#$%&_ '
            texto = pytesseract.image_to_string(image, lang="spa", config=custom_config)

            if len(texto.strip()) < 10:
                logger.warning("OCR devolvió muy poco texto, posible imagen de baja calidad o texto ilegible")
                return "[ADVERTENCIA: Imagen procesada con OCR - texto mínimo extraído]\n" + texto.strip()

            total_chars = len(texto.replace(' ', '').replace('\n', ''))
            weird_chars = sum(1 for char in texto if not char.isalnum() and char not in 'áéíóúüñÁÉÍÓÚÜÑ.,;:!?¿¡()[]{}"-+*/=@#$%&_\n ')

            if total_chars > 0 and weird_chars / total_chars > 0.3:
                logger.warning("OCR detectó muchos caracteres extraños, posible baja calidad de imagen")
                return "[ADVERTENCIA: OCR de baja calidad detectado]\n" + texto.strip()

            return texto.strip()

    except Exception as e:
        logger.error(f"Error procesando imagen {filepath}: {e}")
        raise ValueError(f"No se pudo procesar la imagen: {e}")

def validar_contenido_extraido(texto, tipo_archivo):

    if not texto or len(texto.strip()) < 10:
        return {
            'valido': False,
            'mensaje': f'Contenido muy escaso extraído del archivo {tipo_archivo}',
            'sugerencia': 'Verificar calidad del archivo o usar OCR si contiene imágenes'
        }

    simbolos_extraños = sum(1 for char in texto if char in '�□▯◯●◆■')
    if simbolos_extraños > 10:
        return {
            'valido': False,
            'mensaje': 'Contenido con símbolos no reconocidos detectado',
            'sugerencia': 'Revisar codificación del archivo o usar herramientas de conversión'
        }

    return {
        'valido': True,
        'mensaje': 'Contenido extraído correctamente',
        'sugerencia': None
    }
